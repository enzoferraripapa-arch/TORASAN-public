"""
md2pptx - Obsidian Markdown -> PowerPoint converter

Usage:
    python scripts/md2pptx.py input.md [-o output.pptx] [-t theme]

Markdown format:
    ---
    theme: earth
    font_jp: 游ゴシック
    font_en: Segoe UI
    ---

    # [title] Main Title
    subtitle text
    > key message
    `v1.0` `2026.03`

    ---

    # [content] Slide Title
    - bullet 1
    - bullet 2
      - nested

    ---

    # [two-col] Comparison
    ## Left Column
    - item 1
    ## Right Column
    - item 1

    ---

    # [grid] Grid Title
    ## Card 1
    description
    ## Card 2
    description

    ---

    # [roadmap] Roadmap Title
    > subtitle line
    ## Phase 1 (2026-2027) | Header Label
    ### 可能
    - ◎可能：内容...
    ### 限界
    - △限界：内容...
    ### 施策 | 下段ヘッダー
    - ◇施策内容...

    ## ギャップ分析
    ### ギャップ要因
    - 法的責任
    ### 解決の方向性
    - 段階的導入

    ---

    # [steps] Progression
    ## L1 Draft
    criteria
    ## L2 Active
    criteria

    ---

    # [hub] Hub Title
    central label
    ## Satellite 1
    - description

    ---

    # [summary] Summary
    > key message
    ## Value 1
    detail

Layouts: title, content, two-col, grid, roadmap, steps, hub, summary
Themes:  earth, corporate, ocean, mono, edia
"""

import math
import os
import re
import sys
import argparse

try:
    import yaml
except ImportError:
    yaml = None

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ============================================================
# Themes
# ============================================================
THEMES = {
    "earth": {
        "primary": (0xE0, 0x7A, 0x5F),
        "secondary": (0x81, 0xB2, 0x9A),
        "accent": (0xF2, 0xCC, 0x8F),
        "dark": (0x3D, 0x40, 0x5B),
        "bg": (0xFA, 0xF7, 0xEA),
        "text": (0x3D, 0x40, 0x5B),
        "white": (0xFF, 0xFF, 0xFF),
        "light1": (0xF0, 0xAA, 0x8F),
        "light2": (0xC8, 0xE0, 0xD2),
        "dark2": (0xC0, 0x5A, 0x3F),
        "dark3": (0x5B, 0x75, 0x53),
    },
    "corporate": {
        "primary": (0x00, 0x66, 0xCC),
        "secondary": (0x00, 0x99, 0x66),
        "accent": (0xFF, 0x99, 0x00),
        "dark": (0x1A, 0x1A, 0x2E),
        "bg": (0xF5, 0xF5, 0xF5),
        "text": (0x33, 0x33, 0x33),
        "white": (0xFF, 0xFF, 0xFF),
        "light1": (0xCC, 0xE5, 0xFF),
        "light2": (0xCC, 0xEB, 0xDB),
        "dark2": (0x00, 0x44, 0x88),
        "dark3": (0x00, 0x66, 0x44),
    },
    "ocean": {
        "primary": (0x00, 0x77, 0xB6),
        "secondary": (0x00, 0xB4, 0xD8),
        "accent": (0x90, 0xE0, 0xEF),
        "dark": (0x03, 0x04, 0x5E),
        "bg": (0xF0, 0xF8, 0xFF),
        "text": (0x1A, 0x1A, 0x2E),
        "white": (0xFF, 0xFF, 0xFF),
        "light1": (0xCA, 0xF0, 0xF8),
        "light2": (0xAD, 0xE8, 0xF4),
        "dark2": (0x02, 0x30, 0x47),
        "dark3": (0x00, 0x96, 0xC7),
    },
    "mono": {
        "primary": (0x2D, 0x2D, 0x2D),
        "secondary": (0x5C, 0x5C, 0x5C),
        "accent": (0xE8, 0x4D, 0x3D),
        "dark": (0x1A, 0x1A, 0x1A),
        "bg": (0xFA, 0xFA, 0xFA),
        "text": (0x2D, 0x2D, 0x2D),
        "white": (0xFF, 0xFF, 0xFF),
        "light1": (0xE0, 0xE0, 0xE0),
        "light2": (0xF0, 0xF0, 0xF0),
        "dark2": (0x1A, 0x1A, 0x1A),
        "dark3": (0x3D, 0x3D, 0x3D),
    },
    "edia": {
        "primary": (0x00, 0x00, 0xFF),
        "secondary": (0x00, 0x20, 0x60),
        "accent": (0xFF, 0xFF, 0x00),
        "dark": (0x00, 0x00, 0x00),
        "bg": (0xFF, 0xFF, 0xFF),
        "text": (0x00, 0x00, 0x00),
        "white": (0xFF, 0xFF, 0xFF),
        "light1": (0xCC, 0xE5, 0xFF),
        "light2": (0xE0, 0xFF, 0xE0),
        "dark2": (0x00, 0x20, 0x60),
        "dark3": (0xFF, 0x00, 0x00),
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(t):
    """Tuple -> RGBColor."""
    return RGBColor(*t)


def _lighten(color_tuple, factor):
    """(r,g,b) tuple -> lighter (r,g,b) tuple. factor: 0=original, 1=white."""
    r, g, b = color_tuple
    return (int(r + (255 - r) * factor),
            int(g + (255 - g) * factor),
            int(b + (255 - b) * factor))


# ============================================================
# Markdown parser
# ============================================================

def parse_yaml_frontmatter(text):
    """Simple YAML-like frontmatter parser (no yaml dependency)."""
    result = {}
    for line in text.strip().split("\n"):
        m = re.match(r"^(\w[\w_]*):\s*(.+)", line)
        if m:
            result[m.group(1).strip()] = m.group(2).strip()
    return result


def parse_markdown(md_text):
    """Parse markdown into config + slide list."""
    config = {}
    body = md_text

    fm = re.match(r"^---\s*\n(.*?)\n---\s*\n", md_text, re.DOTALL)
    if fm:
        raw = fm.group(1)
        if yaml:
            config = yaml.safe_load(raw) or {}
        else:
            config = parse_yaml_frontmatter(raw)
        body = md_text[fm.end():]

    blocks = re.split(r"\n---\s*\n", body)
    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        s = parse_slide_block(block)
        if s:
            slides.append(s)
    return config, slides


def parse_slide_block(block):
    """Parse one slide block into structured data."""
    lines = block.split("\n")

    # H1 with optional [layout]
    m = re.match(r"^#\s+\[(\w[\w-]*)\]\s*(.*)", lines[0])
    if m:
        layout = m.group(1)
        title = m.group(2).strip()
    else:
        m2 = re.match(r"^#\s+(.*)", lines[0])
        if m2:
            layout = "content"
            title = m2.group(1).strip()
        else:
            return None

    rest = "\n".join(lines[1:]).strip()

    # Split into H2 sections
    parts = re.split(r"^##\s+", rest, flags=re.MULTILINE)
    preamble = parts[0].strip()

    sections = []
    for part in parts[1:]:
        sec_lines = part.split("\n", 1)
        sec_title = sec_lines[0].strip()
        sec_body = sec_lines[1].strip() if len(sec_lines) > 1 else ""
        sections.append({"title": sec_title, "body": sec_body})

    return {
        "layout": layout,
        "title": title,
        "preamble": preamble,
        "sections": sections,
    }


def parse_bullets(text):
    """Parse bullet list -> [(indent, text), ...]."""
    bullets = []
    for line in text.split("\n"):
        m = re.match(r"^(\s*)[-*]\s+(.*)", line)
        if m:
            indent = len(m.group(1)) // 2
            bullets.append((indent, m.group(2).strip()))
    return bullets


def parse_blockquotes(text):
    """Extract > lines."""
    return [m.group(1) for m in re.finditer(r"^>\s*(.*)", text, re.MULTILINE)]


def extract_tags(text):
    """Extract `tag` backtick tokens."""
    return re.findall(r"`([^`]+)`", text)


def plain_lines(text):
    """Non-bullet, non-quote, non-tag lines."""
    result = []
    for line in text.split("\n"):
        line = line.strip()
        if line and not line.startswith(">") and not line.startswith("-") and "`" not in line:
            result.append(line)
    return result


# ============================================================
# Slide renderer
# ============================================================

class Renderer:
    def __init__(self, prs, theme_name="earth", font_jp="游ゴシック", font_en="Segoe UI"):
        self.prs = prs
        self.theme_name = theme_name
        theme_data = THEMES.get(theme_name, THEMES["earth"])
        self.t = {k: rgb(v) for k, v in theme_data.items()}
        self.t_raw = theme_data
        self.font_jp = font_jp
        self.font_en = font_en
        self.slide_num = 0
        self.total = 0
        self.palette = [
            self.t["primary"], self.t["secondary"], self.t["dark3"],
            self.t["accent"], self.t["dark2"], self.t["light1"],
        ]

    # -- primitives --

    def new_slide(self):
        self.slide_num += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.t["bg"]
        return slide

    def shape(self, slide, stype, l, t, w, h, fill=None, line=None, line_w=Pt(2)):
        s = slide.shapes.add_shape(stype, l, t, w, h)
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        else:
            s.fill.background()
        if line:
            s.line.color.rgb = line
            s.line.width = line_w
        else:
            s.line.fill.background()
        return s

    def tbox(self, slide, l, t, w, h, text, sz=Pt(14), color=None,
             bold=False, align=PP_ALIGN.LEFT, font=None):
        color = color or self.t["text"]
        font = font or self.font_jp
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = sz
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        return tf

    def rrect(self, slide, l, t, w, h, fill, text="", sz=Pt(12),
              tc=None, bold=False, line=None):
        tc = tc or self.t["white"]
        s = self.shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
                       fill=fill, line=line)
        if text:
            tf = s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = sz
            p.font.color.rgb = tc
            p.font.bold = bold
            p.font.name = self.font_jp
            p.alignment = PP_ALIGN.CENTER
        return s

    def header(self, slide, title):
        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(0), SLIDE_W, Inches(1.1),
                   fill=self.t["dark"])
        self.tbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                  title, sz=Pt(30), color=self.t["white"], bold=True)

    def footer(self, slide):
        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(7.2), SLIDE_W, Inches(0.3),
                   fill=self.t["secondary"])
        self.tbox(slide, Inches(12.3), Inches(7.05), Inches(1), Inches(0.4),
                  f"{self.slide_num} / {self.total}",
                  sz=Pt(10), color=self.t["secondary"], align=PP_ALIGN.RIGHT)

    def add_para(self, tf, text, sz=Pt(14), color=None, bold=False,
                 align=PP_ALIGN.LEFT, space_before=Pt(4)):
        color = color or self.t["text"]
        p = tf.add_paragraph()
        p.text = text
        p.font.size = sz
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = self.font_jp
        p.alignment = align
        p.space_before = space_before
        return p

    def write_bullets(self, tf, bullets, base_sz=14):
        for i, (indent, text) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = "    " * indent + ("- " if indent > 0 else "")
            p.text = prefix + text
            p.font.size = Pt(max(9, base_sz - indent * 2))
            p.font.color.rgb = self.t["text"]
            p.font.name = self.font_jp
            p.space_before = Pt(4)

    # -- decorations (swoosh / shadow) --

    def _add_swoosh(self, slide, left, top, width, height, rot_deg, lower=False):
        """Add swooshArrow background with gradient. Coordinates in inches."""
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
        spPr = s._element.spPr

        # Rotation
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.set("rot", str(int(rot_deg * 60000)))

        # Change geometry to swooshArrow
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "swooshArrow")
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = prstGeom.makeelement(qn("a:avLst"), {})
                prstGeom.append(avLst)
            else:
                avLst.clear()
            avLst.append(avLst.makeelement(
                qn("a:gd"), {"name": "adj1", "fmla": "val 25000"}))
            avLst.append(avLst.makeelement(
                qn("a:gd"), {"name": "adj2", "fmla": "val 25000"}))

        # Gradient fill derived from theme colors
        base = self.t_raw["secondary"] if lower else self.t_raw["primary"]
        c1 = _lighten(base, 0.55)
        c2 = _lighten(base, 0.85)

        for tag in ["a:solidFill", "a:gradFill", "a:noFill"]:
            old = spPr.find(qn(tag))
            if old is not None:
                spPr.remove(old)

        gradFill = spPr.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
        gsLst = gradFill.makeelement(qn("a:gsLst"), {})

        gs1 = gsLst.makeelement(qn("a:gs"), {"pos": "0"})
        clr1 = gs1.makeelement(qn("a:srgbClr"),
                               {"val": "%02X%02X%02X" % c1})
        clr1.append(clr1.makeelement(qn("a:alpha"), {"val": "80000"}))
        gs1.append(clr1)
        gsLst.append(gs1)

        gs2 = gsLst.makeelement(qn("a:gs"), {"pos": "100000"})
        clr2 = gs2.makeelement(qn("a:srgbClr"),
                               {"val": "%02X%02X%02X" % c2})
        clr2.append(clr2.makeelement(qn("a:alpha"), {"val": "50000"}))
        gs2.append(clr2)
        gsLst.append(gs2)

        gradFill.append(gsLst)
        gradFill.append(gradFill.makeelement(
            qn("a:lin"), {"ang": "10800000", "scaled": "1"}))

        prstGeom_el = spPr.find(qn("a:prstGeom"))
        if prstGeom_el is not None:
            prstGeom_el.addnext(gradFill)
        else:
            spPr.append(gradFill)

        # No line
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = spPr.makeelement(qn("a:ln"), {})
            spPr.append(ln)
        else:
            ln.clear()
        ln.append(ln.makeelement(qn("a:noFill"), {}))

        # Glow effect
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        glow = effectLst.makeelement(qn("a:glow"), {"rad": "127000"})
        srgb = glow.makeelement(qn("a:srgbClr"), {"val": "4BACC6"})
        srgb.append(srgb.makeelement(qn("a:lumMod"), {"val": "20000"}))
        srgb.append(srgb.makeelement(qn("a:lumOff"), {"val": "80000"}))
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": "10000"}))
        glow.append(srgb)
        effectLst.append(glow)
        spPr.append(effectLst)

        return s

    def _add_shadow(self, shape):
        """Add drop shadow effect."""
        spPr = shape._element.spPr
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
            "blurRad": "40000", "dist": "23000",
            "dir": "5400000", "rotWithShape": "0"
        })
        srgb = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": "38000"}))
        outerShdw.append(srgb)
        effectLst.append(outerShdw)
        spPr.append(effectLst)

    def _render_gap_callout(self, slide, gap_sec, x=0.12, y=5.98,
                            w=2.94, h=0.92):
        """Render gap analysis section as callout box."""
        t = self.t
        s = self.shape(slide, MSO_SHAPE.RECTANGULAR_CALLOUT,
                       Inches(x), Inches(y), Inches(w), Inches(h),
                       fill=t["secondary"])
        # Change to wedgeRectCallout with pointer up
        prstGeom = s._element.spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "wedgeRectCallout")
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = prstGeom.makeelement(qn("a:avLst"), {})
                prstGeom.append(avLst)
            else:
                avLst.clear()
            avLst.append(avLst.makeelement(
                qn("a:gd"), {"name": "adj1", "fmla": "val 60180"}))
            avLst.append(avLst.makeelement(
                qn("a:gd"), {"name": "adj2", "fmla": "val -95818"}))
        self._add_shadow(s)

        # Render H3 subsections as columns inside the callout
        body = gap_sec["body"]
        h3_parts = re.split(r"^###\s+", body, flags=re.MULTILINE)

        col_x = x
        col_w = w / max(len(h3_parts) - 1, 1)
        for h3 in h3_parts[1:]:
            h3_lines = h3.split("\n", 1)
            h3_title = h3_lines[0].strip()
            h3_body = h3_lines[1].strip() if len(h3_lines) > 1 else ""
            bullets = [bt for _, bt in parse_bullets(h3_body)]

            tf = self.tbox(slide, Inches(col_x), Inches(y - 0.37),
                          Inches(col_w), Inches(h + 0.7),
                          h3_title, sz=Pt(9), color=t["white"], bold=True)
            for bt in bullets:
                self.add_para(tf, f"◇{bt}", sz=Pt(9), color=t["white"],
                             bold=True, space_before=Pt(0))
            col_x += col_w

    def _auto_color_bullet(self, text):
        """Return appropriate color for a bullet based on its prefix."""
        t = self.t
        if text.startswith("◎") or text.startswith("可能"):
            return t["primary"]
        elif text.startswith("△") or text.startswith("▲") or text.startswith("限界"):
            return t["dark3"]
        elif text.startswith("◇") or text.startswith("◆"):
            return t["secondary"]
        elif text.startswith("→"):
            return t["secondary"]
        return t["text"]

    # -- layouts --

    def lay_title(self, d):
        slide = self.new_slide()
        t = self.t

        # Decorative circles
        for cx, cy, diam, c in [
            (10.0, 4.0, 5.0, "light1"), (10.8, 4.8, 4.0, "accent"),
            (-1.0, -1.0, 3.0, "light2"), (-0.3, 5.5, 2.5, "light2"),
        ]:
            self.shape(slide, MSO_SHAPE.OVAL,
                       Inches(cx), Inches(cy), Inches(diam), Inches(diam),
                       fill=t[c])

        # Accent bar
        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.5),
                   fill=t["primary"])

        # Title
        self.tbox(slide, Inches(1.3), Inches(1.8), Inches(8), Inches(1.2),
                  d["title"], sz=Pt(48), color=t["dark"], bold=True)

        # Subtitle / quote / tags from preamble
        y = Inches(3.1)
        for line in plain_lines(d["preamble"]):
            self.tbox(slide, Inches(1.3), y, Inches(8), Inches(0.7),
                      line, sz=Pt(24), color=t["primary"], bold=True)
            y += Inches(0.9)
            break  # first line only

        for q in parse_blockquotes(d["preamble"]):
            self.tbox(slide, Inches(1.3), y, Inches(8), Inches(0.7),
                      q, sz=Pt(16), color=t["text"])
            y += Inches(0.8)
            break

        tags = extract_tags(d["preamble"])
        if tags:
            x = Inches(1.3)
            for tag in tags:
                w = Inches(max(1.5, len(tag) * 0.18))
                self.rrect(slide, x, y, w, Inches(0.5),
                           t["secondary"], tag, sz=Pt(14), bold=True)
                x += w + Inches(0.3)

        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
                   fill=t["primary"])
        self.tbox(slide, Inches(12.3), Inches(7.05), Inches(1), Inches(0.4),
                  f"{self.slide_num} / {self.total}",
                  sz=Pt(10), color=t["secondary"], align=PP_ALIGN.RIGHT)

    def lay_content(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])

        y = Inches(1.5)

        # Preamble bullets
        bullets = parse_bullets(d["preamble"])
        if bullets:
            tf = self.tbox(slide, Inches(0.8), y, Inches(11.5), Inches(5), "")
            self.write_bullets(tf, bullets)
            y += Inches(len(bullets) * 0.42)

        # Sections
        for sec in d["sections"]:
            self.tbox(slide, Inches(0.8), y, Inches(11.5), Inches(0.5),
                      sec["title"], sz=Pt(18), color=self.t["primary"], bold=True)
            y += Inches(0.55)
            sb = parse_bullets(sec["body"])
            if sb:
                tf = self.tbox(slide, Inches(1.2), y, Inches(11), Inches(3), "")
                self.write_bullets(tf, sb)
                y += Inches(len(sb) * 0.4)
            y += Inches(0.15)

        self.footer(slide)

    def lay_two_col(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])
        colors = [self.t["primary"], self.t["secondary"]]

        for i, sec in enumerate(d["sections"][:2]):
            cx = Inches(0.5 + i * 6.3)
            self.rrect(slide, cx, Inches(1.4), Inches(5.8), Inches(0.7),
                       colors[i], sec["title"], sz=Pt(20), bold=True)

            bullets = parse_bullets(sec["body"])
            if bullets:
                tf = self.tbox(slide, cx + Inches(0.2), Inches(2.3),
                               Inches(5.4), Inches(4.5), "")
                self.write_bullets(tf, bullets)
            else:
                lines = [l.strip() for l in sec["body"].split("\n") if l.strip()]
                if lines:
                    tf = self.tbox(slide, cx + Inches(0.2), Inches(2.3),
                                   Inches(5.4), Inches(4.5), lines[0])
                    for line in lines[1:]:
                        self.add_para(tf, line)

        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(6.55), Inches(1.4), Inches(0.05), Inches(5.3),
                   fill=self.t["accent"])
        self.footer(slide)

    def lay_grid(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])
        secs = d["sections"]
        n = len(secs)
        cols = min(n, 4) if n > 3 else n
        rows = (n + cols - 1) // cols

        cw = Inches(min(3.8, 12.0 / cols))
        ch = Inches(min(2.5, 5.5 / rows))
        gap = Inches(0.3)
        total_w = int(cols * cw + (cols - 1) * gap)
        sx = (int(SLIDE_W) - total_w) // 2

        for i, sec in enumerate(secs):
            c, r = i % cols, i // cols
            x = sx + c * int(cw + gap)
            y = int(Inches(1.4)) + r * int(ch + gap)
            color = self.palette[i % len(self.palette)]

            s = self.rrect(slide, x, y, cw, ch, color)
            tf = s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sec["title"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.t["white"]
            p.font.name = self.font_jp
            p.alignment = PP_ALIGN.CENTER

            for line in sec["body"].split("\n"):
                line = re.sub(r"^\s*[-*]\s+", "", line).strip()
                if not line:
                    continue
                self.add_para(tf, line, sz=Pt(11),
                              color=RGBColor(0xFF, 0xFF, 0xDD),
                              align=PP_ALIGN.CENTER, space_before=Pt(3))

        self.footer(slide)

    def lay_roadmap(self, d):
        """Enhanced roadmap with auto-layout, swoosh backgrounds, and diagonal flow."""
        slide = self.new_slide()
        secs = d["sections"]
        t = self.t

        # Separate phases from special sections (ギャップ分析 etc.)
        phases = []
        gap_sec = None
        for sec in secs:
            if "ギャップ" in sec["title"] or "gap" in sec["title"].lower():
                gap_sec = sec
            else:
                phases.append(sec)

        n = len(phases)
        if n == 0:
            self.footer(slide)
            return

        swoosh = self.theme_name in ("edia",)
        diagonal = n >= 3
        quotes = parse_blockquotes(d["preamble"])

        # ---- Title area ----
        if swoosh:
            self.tbox(slide, Inches(0.16), Inches(0.06),
                      Inches(13), Inches(0.54),
                      d["title"], sz=Pt(28), color=t["text"], bold=True)
            if quotes:
                bar = self.shape(slide, MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0.74), SLIDE_W, Inches(0.69),
                                 fill=t["accent"], line=t["dark3"], line_w=Pt(3))
                tf = self.tbox(slide, Inches(0.1), Inches(0.76),
                              Inches(13.1), Inches(0.64),
                              quotes[0], sz=Pt(18), color=t["primary"],
                              bold=True, align=PP_ALIGN.CENTER)
                if len(quotes) > 1:
                    self.add_para(tf, quotes[1], sz=Pt(18), color=t["text"],
                                 bold=True, align=PP_ALIGN.CENTER)
            content_top = 1.5
        else:
            self.header(slide, d["title"])
            if quotes:
                self.tbox(slide, Inches(0.5), Inches(1.15),
                          Inches(12), Inches(0.35),
                          quotes[0], sz=Pt(12), color=t["text"],
                          align=PP_ALIGN.CENTER)
            content_top = 1.65

        # ---- Swoosh backgrounds ----
        if swoosh:
            self._add_swoosh(slide, 0, 2.41, 13.0, 4.0, 355.0, lower=False)
            self._add_swoosh(slide, 0.2, 4.38, 12.5, 3.3, 1.5, lower=True)

        # ---- Legend (edia-style) ----
        if swoosh:
            # Small legend in top-left
            lg = self.shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.09), Inches(content_top + 0.12),
                            Inches(3.01), Inches(0.46),
                            fill=RGBColor(0xF5, 0xF5, 0xF5))
            lg.line.fill.background()
            ltf = self.tbox(slide, Inches(0.12), Inches(content_top + 0.12),
                           Inches(2.95), Inches(0.15),
                           "【凡例】", sz=Pt(9), color=t["text"], bold=True)
            self.add_para(ltf, "◎可能：技術的に実用可能な領域",
                         sz=Pt(9), color=t["text"], bold=False, space_before=Pt(0))
            self.add_para(ltf, "△限界：法規制・倫理・情緒的要因で制約",
                         sz=Pt(9), color=t["text"], bold=False, space_before=Pt(0))

        # ---- Auto-layout constants ----
        margin = 0.1
        usable = 13.333 - 2 * margin
        col_w = usable / n
        hdr_h = 0.55

        # Diagonal range: Phase 1 bottom-left → Phase N top-right
        if diagonal:
            y_range = min(2.0, 0.55 * (n - 1))
            y_base = content_top + y_range + 0.5
        else:
            y_range = 0
            y_base = content_top + 0.5

        has_lower = False
        for phase in phases:
            if re.search(r"^###\s+.*施策", phase["body"], re.MULTILINE):
                has_lower = True
                break

        for i, phase in enumerate(phases):
            x = margin + col_w * i
            if diagonal and n > 1:
                phase_y = y_base - y_range * i / (n - 1)
            else:
                phase_y = y_base

            # Parse title: "Phase 1 (2026-2027) | Header Label"
            parts = phase["title"].split("|")
            ptitle = parts[0].strip()
            plabel = parts[1].strip() if len(parts) > 1 else ""

            # Parse H3 sections into upper/lower
            body = phase["body"]
            h3_parts = re.split(r"^###\s+", body, flags=re.MULTILINE)
            preamble_bullets = [bt for _, bt in parse_bullets(h3_parts[0])]

            upper_bullets = []
            lower_header = ""
            lower_bullets = []

            for h3 in h3_parts[1:]:
                h3_lines = h3.split("\n", 1)
                h3_title = h3_lines[0].strip()
                h3_body = h3_lines[1].strip() if len(h3_lines) > 1 else ""
                bullets = [bt for _, bt in parse_bullets(h3_body)]

                if "施策" in h3_title or "組織" in h3_title:
                    if "|" in h3_title:
                        lower_header = h3_title.split("|", 1)[1].strip()
                    else:
                        lh = h3_title.replace("施策", "").strip()
                        lower_header = lh.lstrip(":").strip() if lh else ""
                    lower_bullets.extend(bullets)
                else:
                    upper_bullets.extend(bullets)

            # Fallback: use preamble bullets if no H3
            if not upper_bullets:
                upper_bullets = preamble_bullets

            # ---- Vertical divider ----
            if i > 0:
                div_y_top = phase_y - 0.3
                div_h = 5.5 if has_lower else 3.0
                cxn = slide.shapes.add_connector(
                    1, Inches(x - 0.02), Inches(div_y_top),
                    Inches(x - 0.02), Inches(div_y_top + div_h))
                cxn.line.width = Pt(2)
                cxn.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)

            # ---- Phase label ----
            self.tbox(slide, Inches(x + 0.05), Inches(phase_y - 0.28),
                      Inches(col_w - 0.1), Inches(0.27),
                      ptitle, sz=Pt(16), color=t["text"], bold=True)

            # ---- Upper header ----
            hdr_text = plabel or ptitle
            hdr = self.rrect(slide, Inches(x), Inches(phase_y),
                             Inches(col_w - 0.1), Inches(hdr_h),
                             t["primary"], hdr_text, sz=Pt(14), bold=True)
            if swoosh:
                self._add_shadow(hdr)

            # ---- Upper content ----
            cy = phase_y + hdr_h + 0.05
            if upper_bullets:
                tf = self.tbox(slide, Inches(x), Inches(cy),
                              Inches(col_w - 0.05), Inches(1.5), "")
                for bi, bt in enumerate(upper_bullets):
                    color = self._auto_color_bullet(bt)
                    if bi == 0:
                        tf.paragraphs[0].text = bt
                        tf.paragraphs[0].font.size = Pt(9)
                        tf.paragraphs[0].font.color.rgb = color
                        tf.paragraphs[0].font.bold = True
                        tf.paragraphs[0].font.name = self.font_jp
                    else:
                        self.add_para(tf, bt, sz=Pt(9), color=color,
                                     bold=True, space_before=Pt(0))

            # ---- Lower section ----
            if lower_bullets or lower_header:
                lower_y = phase_y + hdr_h + 1.7

                # Lower phase label
                self.tbox(slide, Inches(x + 0.05), Inches(lower_y - 0.28),
                          Inches(col_w - 0.1), Inches(0.27),
                          ptitle, sz=Pt(16), color=t["text"], bold=True)

                # Lower header
                if lower_header:
                    lhdr = self.rrect(slide, Inches(x), Inches(lower_y),
                                      Inches(col_w - 0.1), Inches(hdr_h),
                                      t["secondary"], lower_header,
                                      sz=Pt(14), bold=True)
                    if swoosh:
                        self._add_shadow(lhdr)

                # Lower content
                lcy = lower_y + hdr_h + 0.05
                if lower_bullets:
                    tf = self.tbox(slide, Inches(x), Inches(lcy),
                                  Inches(col_w - 0.05), Inches(1.8), "")
                    for bi, bt in enumerate(lower_bullets):
                        color = self._auto_color_bullet(bt)
                        if bi == 0:
                            tf.paragraphs[0].text = bt
                            tf.paragraphs[0].font.size = Pt(9)
                            tf.paragraphs[0].font.color.rgb = color
                            tf.paragraphs[0].font.bold = True
                            tf.paragraphs[0].font.name = self.font_jp
                        else:
                            self.add_para(tf, bt, sz=Pt(9), color=color,
                                         bold=True, space_before=Pt(0))

        # ---- Gap analysis callout ----
        if gap_sec:
            self._render_gap_callout(slide, gap_sec)

        # ---- Bottom-right legend badges ----
        if swoosh:
            self.rrect(slide, Inches(11.88), Inches(6.43),
                       Inches(1.42), Inches(0.35),
                       t["primary"], "技術的可能性分析", sz=Pt(9), bold=True)
            self.rrect(slide, Inches(11.88), Inches(6.83),
                       Inches(1.42), Inches(0.35),
                       t["secondary"], "組織施策", sz=Pt(9), bold=True)

        if not swoosh:
            # Simple timeline arrow for non-swoosh themes
            self.shape(slide, MSO_SHAPE.RIGHT_ARROW,
                       Inches(margin), Inches(6.7),
                       Inches(usable), Inches(0.25),
                       fill=t["accent"])
            self.footer(slide)

    def lay_steps(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])
        secs = d["sections"]
        n = len(secs)
        palette = [self.t["light2"], self.t["secondary"], self.t["accent"],
                   self.t["primary"], self.t["dark2"]]

        for i, sec in enumerate(secs):
            x = Inches(0.5 + i * 1.0)
            y = Inches(5.5 - i * (4.0 / max(n - 1, 1)))
            w = Inches(12.0 - i * 1.0)

            s = self.rrect(slide, x, y, w, Inches(0.7),
                           palette[i % len(palette)])
            tf = s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sec["title"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.t["white"]
            p.font.name = self.font_en

            if sec["body"]:
                first = re.sub(r"^\s*[-*]\s+", "", sec["body"].split("\n")[0]).strip()
                if first:
                    self.add_para(tf, first, sz=Pt(11),
                                  color=RGBColor(0xFF, 0xFF, 0xDD))

        self.footer(slide)

    def lay_hub(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])
        secs = d["sections"]
        hub_text = plain_lines(d["preamble"])[0] if plain_lines(d["preamble"]) else d["title"]

        # Hub circle
        hub = self.shape(slide, MSO_SHAPE.OVAL,
                         Inches(5.0), Inches(2.3), Inches(3.2), Inches(3.2),
                         fill=self.t["primary"])
        tf = hub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = hub_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.t["white"]
        p.alignment = PP_ALIGN.CENTER

        n = len(secs)
        for i, sec in enumerate(secs):
            angle = (2 * math.pi * i / n) - math.pi / 2
            cx = int(Inches(6.6) + Inches(4.2) * math.cos(angle))
            cy = int(Inches(3.9) + Inches(2.8) * math.sin(angle))
            color = self.palette[i % len(self.palette)]

            box = self.rrect(slide, cx - int(Inches(1.5)), cy - int(Inches(0.5)),
                             Inches(3.0), Inches(1.0), color)
            tf2 = box.text_frame
            tf2.word_wrap = True
            p = tf2.paragraphs[0]
            p.text = sec["title"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.t["white"]
            p.alignment = PP_ALIGN.CENTER

            if sec["body"]:
                first = re.sub(r"^\s*[-*]\s+", "", sec["body"].split("\n")[0]).strip()
                if first:
                    self.add_para(tf2, first, sz=Pt(10),
                                  color=RGBColor(0xFF, 0xFF, 0xDD),
                                  align=PP_ALIGN.CENTER)

        self.footer(slide)

    def lay_summary(self, d):
        slide = self.new_slide()
        self.header(slide, d["title"])
        secs = d["sections"]
        n = len(secs)

        if n > 0:
            cw = Inches(min(3.5, 11.0 / n))
            total_w = int(n * cw)
            sx = (int(SLIDE_W) - total_w) // 2

            for i, sec in enumerate(secs):
                x = sx + i * int(cw)
                color = self.palette[i % len(self.palette)]
                diam = int(cw) - int(Inches(0.3))

                s = self.shape(slide, MSO_SHAPE.OVAL,
                               x, Inches(1.5), diam, diam, fill=color)
                tf = s.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sec["title"]
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = self.t["white"]
                p.alignment = PP_ALIGN.CENTER

                for line in sec["body"].split("\n")[:3]:
                    line = re.sub(r"^\s*[-*]\s+", "", line).strip()
                    if line:
                        self.add_para(tf, line, sz=Pt(10),
                                      color=RGBColor(0xFF, 0xFF, 0xEE),
                                      align=PP_ALIGN.CENTER, space_before=Pt(2))

        # Key message
        quotes = parse_blockquotes(d["preamble"])
        if quotes:
            self.shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(1), Inches(5.5), Inches(11.3), Inches(1.2),
                       fill=self.t["white"], line=self.t["primary"])
            self.tbox(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.8),
                      quotes[0], sz=Pt(18), color=self.t["text"],
                      bold=True, align=PP_ALIGN.CENTER)

        self.shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
                   fill=self.t["primary"])
        self.tbox(slide, Inches(12.3), Inches(7.05), Inches(1), Inches(0.4),
                  f"{self.slide_num} / {self.total}",
                  sz=Pt(10), color=self.t["secondary"], align=PP_ALIGN.RIGHT)

    # -- dispatch --

    LAYOUTS = {
        "title": lay_title,
        "content": lay_content,
        "bullets": lay_content,
        "two-col": lay_two_col,
        "grid": lay_grid,
        "roadmap": lay_roadmap,
        "steps": lay_steps,
        "hub": lay_hub,
        "summary": lay_summary,
    }

    def render_all(self, slides):
        self.total = len(slides)
        for d in slides:
            fn = self.LAYOUTS.get(d["layout"], self.lay_content)
            fn(self, d)


# ============================================================
# Public API
# ============================================================

def convert(input_path, output_path=None, theme_override=None):
    """Convert markdown file to pptx."""
    with open(input_path, "r", encoding="utf-8") as f:
        md_text = f.read()

    config, slides = parse_markdown(md_text)
    theme = theme_override or config.get("theme", "earth")
    font_jp = config.get("font_jp", "游ゴシック")
    font_en = config.get("font_en", "Segoe UI")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    renderer = Renderer(prs, theme, font_jp, font_en)
    renderer.render_all(slides)

    if not output_path:
        output_path = os.path.splitext(input_path)[0] + ".pptx"

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"Generated: {output_path} ({len(slides)} slides)")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Markdown -> PowerPoint converter")
    parser.add_argument("input", help="Input markdown file")
    parser.add_argument("-o", "--output", help="Output pptx path")
    parser.add_argument("-t", "--theme", choices=list(THEMES.keys()),
                        help="Theme: earth, corporate, ocean, mono, edia")
    args = parser.parse_args()
    convert(args.input, args.output, args.theme)


if __name__ == "__main__":
    main()
