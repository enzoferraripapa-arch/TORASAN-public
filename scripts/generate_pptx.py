"""
TORASAN スキル & ナレッジ フレームワーク PowerPoint 生成スクリプト
明るいアースカラー調 / クールポップデザイン
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# カラーパレット（明るいアースカラー）
# ============================================================
TERRACOTTA   = RGBColor(0xE0, 0x7A, 0x5F)  # メインアクセント
SAGE         = RGBColor(0x81, 0xB2, 0x9A)  # サブアクセント
SAND         = RGBColor(0xF4, 0xF1, 0xDE)  # ベースカラー
CHARCOAL     = RGBColor(0x3D, 0x40, 0x5B)  # テキスト
MUSTARD      = RGBColor(0xF2, 0xCC, 0x8F)  # ハイライト
OLIVE        = RGBColor(0x5B, 0x75, 0x53)  # ダークアクセント
MINT         = RGBColor(0xB5, 0xE2, 0xD5)  # ライトサブ
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TERRA  = RGBColor(0xF0, 0xAA, 0x8F)  # 薄いテラコッタ
LIGHT_SAGE   = RGBColor(0xC8, 0xE0, 0xD2)  # 薄いセージ
DEEP_TERRA   = RGBColor(0xC0, 0x5A, 0x3F)  # 濃いテラコッタ
SOFT_SAND    = RGBColor(0xFA, 0xF7, 0xEA)  # 薄いサンド

# フォント
FONT_JP = "游ゴシック"
FONT_EN = "Segoe UI"

# スライドサイズ: ワイドスクリーン (13.33 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def cm(val):
    """cm → Emu 変換"""
    return Emu(int(val * 360000))


def set_slide_bg(slide, color):
    """スライド背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=Pt(0)):
    """図形を追加"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=Pt(14), color=CHARCOAL,
             bold=False, alignment=PP_ALIGN.CENTER, font_name=FONT_JP):
    """図形にテキストを設定"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_text_box(slide, left, top, width, height, text,
                 font_size=Pt(14), color=CHARCOAL, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name=FONT_JP):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=Pt(14), color=CHARCOAL,
                  bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_JP,
                  space_before=Pt(4), space_after=Pt(2)):
    """テキストフレームに段落を追加"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_arrow(slide, left, top, width, height, fill_color=TERRACOTTA,
              rotation=0):
    """矢印図形を追加"""
    shape = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
                      fill_color=fill_color)
    shape.rotation = rotation
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=Pt(12), text_color=WHITE,
                     bold=False, line_color=None):
    """角丸矩形 + テキスト"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      left, top, width, height,
                      fill_color=fill_color, line_color=line_color)
    if text:
        set_text(shape, text, font_size=font_size, color=text_color, bold=bold)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.auto_size = None

    # 垂直中央揃え
    shape.text_frame.word_wrap = True
    return shape


def add_circle(slide, left, top, diameter, fill_color,
               text="", font_size=Pt(14), text_color=WHITE, bold=False):
    """円形 + テキスト"""
    shape = add_shape(slide, MSO_SHAPE.OVAL,
                      left, top, diameter, diameter,
                      fill_color=fill_color)
    if text:
        tf = set_text(shape, text, font_size=font_size, color=text_color,
                      bold=bold)
    return shape


def add_deco_circles(slide):
    """装飾用の半透明風の円を追加"""
    add_circle(slide, Inches(11.5), Inches(5.8), Inches(2.5), LIGHT_TERRA)
    add_circle(slide, Inches(12.0), Inches(6.3), Inches(1.8), MUSTARD)
    add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(1.5), MINT)


def add_slide_number(slide, num, total=10):
    """スライド番号"""
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(1.0), Inches(0.4),
                 f"{num} / {total}", font_size=Pt(10), color=SAGE,
                 alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_bottom_bar(slide):
    """ボトムバー (アクセントライン)"""
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(7.2), SLIDE_W, Inches(0.3),
              fill_color=SAGE)


# ============================================================
# Slide 1: タイトル
# ============================================================
def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, SAND)

    # 装飾円
    add_circle(slide, Inches(10.0), Inches(4.0), Inches(5.0), LIGHT_TERRA)
    add_circle(slide, Inches(10.8), Inches(4.8), Inches(4.0), MUSTARD)
    add_circle(slide, Inches(-1.0), Inches(-1.0), Inches(3.0), MINT)
    add_circle(slide, Inches(-0.3), Inches(5.5), Inches(2.5), LIGHT_SAGE)

    # 左サイドのアクセントバー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.5),
              fill_color=TERRACOTTA)

    # タイトル
    add_text_box(slide, Inches(1.3), Inches(1.8), Inches(8.0), Inches(1.2),
                 "TORASAN", font_size=Pt(54), color=CHARCOAL,
                 bold=True, font_name=FONT_EN)

    add_text_box(slide, Inches(1.3), Inches(2.9), Inches(8.0), Inches(0.9),
                 "スキル & ナレッジ フレームワーク",
                 font_size=Pt(32), color=TERRACOTTA, bold=True)

    # サブタイトル
    add_text_box(slide, Inches(1.3), Inches(4.0), Inches(8.0), Inches(0.7),
                 "AI駆動の機能安全開発を支える知識管理システム",
                 font_size=Pt(18), color=CHARCOAL)

    # バージョン & 日付
    add_rounded_rect(slide, Inches(1.3), Inches(5.0), Inches(1.8), Inches(0.5),
                     SAGE, "v4.0", font_size=Pt(16), text_color=WHITE, bold=True)

    add_text_box(slide, Inches(3.3), Inches(5.05), Inches(3.0), Inches(0.5),
                 "2026.03", font_size=Pt(16), color=SAGE,
                 font_name=FONT_EN)

    # ボトムバー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
              fill_color=TERRACOTTA)

    add_slide_number(slide, 1)


# ============================================================
# Slide 2: TORASANとは？
# ============================================================
def slide_02_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)
    add_deco_circles(slide)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "TORASANとは？",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 中央ハブ円
    hub = add_circle(slide, Inches(5.2), Inches(2.2), Inches(2.8), TERRACOTTA,
                     "TORASAN\nHub", font_size=Pt(22), text_color=WHITE, bold=True)

    # 4つの衛星ボックス
    satellites = [
        (Inches(1.0), Inches(2.0), "41 スキル", "実行する知恵", SAGE),
        (Inches(9.5), Inches(2.0), "16 ナレッジ", "参照する知識", OLIVE),
        (Inches(1.0), Inches(4.8), "管理 GUI", "localhost:3000", MUSTARD),
        (Inches(9.5), Inches(4.8), "配布機構", "install.sh & sync", DEEP_TERRA),
    ]

    for x, y, title, sub, color in satellites:
        box = add_rounded_rect(slide, x, y, Inches(3.0), Inches(1.5), color,
                               bold=True, font_size=Pt(18), text_color=WHITE)
        tf = box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xDD)
        p2.alignment = PP_ALIGN.CENTER

    # 矢印（ハブから各衛星へ）
    add_arrow(slide, Inches(4.2), Inches(3.1), Inches(1.0), Inches(0.4),
              LIGHT_TERRA, rotation=180)  # 左
    add_arrow(slide, Inches(8.1), Inches(3.1), Inches(1.0), Inches(0.4),
              LIGHT_TERRA)  # 右
    add_arrow(slide, Inches(4.2), Inches(4.5), Inches(1.0), Inches(0.4),
              LIGHT_TERRA, rotation=180)  # 左下
    add_arrow(slide, Inches(8.1), Inches(4.5), Inches(1.0), Inches(0.4),
              LIGHT_TERRA)  # 右下

    # 説明テキスト
    add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.6),
                 "全スキルの開発・メンテナンスの拠点。"
                 "汎用スキルはinstall.shで全PJに配布、ドメインスキルはsyncで個別配布。",
                 font_size=Pt(14), color=CHARCOAL)

    add_bottom_bar(slide)
    add_slide_number(slide, 2)


# ============================================================
# Slide 3: スキルとナレッジの違い
# ============================================================
def slide_03_skill_vs_knowledge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "スキル vs ナレッジ",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 左: スキル
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.8),
                     TERRACOTTA, "スキル = 実行する知恵",
                     font_size=Pt(22), text_color=WHITE, bold=True)

    skill_items = [
        "コマンドとして呼び出す（例: /session, /driver-gen）",
        "ステップバイステップの手順を定義",
        "入力 → 処理 → 出力 の実行フロー",
        "他スキルと連携してワークフローを構成",
        "成熟度モデル（L1→L5）で品質を管理",
    ]
    tf = add_text_box(slide, Inches(1.0), Inches(2.5), Inches(5.0), Inches(3.5),
                      "", font_size=Pt(14), color=CHARCOAL)
    tf.paragraphs[0].text = ""
    for i, item in enumerate(skill_items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.font.name = FONT_JP
        p.space_before = Pt(6)

    # 右: ナレッジ
    add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.8),
                     SAGE, "ナレッジ = 参照する知識",
                     font_size=Pt(22), text_color=WHITE, bold=True)

    know_items = [
        "スキル実行時に参照されるリファレンス",
        "規格要件・技術仕様・ベストプラクティス",
        ".claude/knowledge/ に Markdown で格納",
        "スキルが自動的に必要なナレッジを読み込む",
        "ドメイン知識を構造化して再利用可能に",
    ]
    tf2 = add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(3.5),
                       "", font_size=Pt(14), color=CHARCOAL)
    tf2.paragraphs[0].text = ""
    for i, item in enumerate(know_items):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = f"  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.font.name = FONT_JP
        p.space_before = Pt(6)

    # 中央: 関係矢印
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(6.3), Inches(1.5), Inches(0.05), Inches(4.5),
              fill_color=MUSTARD)

    # 下部: 関係図
    add_rounded_rect(slide, Inches(2.5), Inches(5.8), Inches(2.5), Inches(0.7),
                     TERRACOTTA, "/driver-gen",
                     font_size=Pt(14), text_color=WHITE, bold=True)
    add_arrow(slide, Inches(5.1), Inches(5.9), Inches(1.5), Inches(0.4),
              MUSTARD)
    add_rounded_rect(slide, Inches(6.7), Inches(5.8), Inches(3.0), Inches(0.7),
                     SAGE, "misra_c_2012.md",
                     font_size=Pt(14), text_color=WHITE, bold=True)

    add_text_box(slide, Inches(5.2), Inches(6.5), Inches(2.0), Inches(0.4),
                 "参照", font_size=Pt(12), color=MUSTARD,
                 alignment=PP_ALIGN.CENTER, bold=True)

    add_bottom_bar(slide)
    add_slide_number(slide, 3)


# ============================================================
# Slide 4: スキルの全体像（41本）
# ============================================================
def slide_04_skill_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "スキルの全体像 — 41本",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 汎用スキル ラベル
    add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(3.0), Inches(0.6),
                     SAGE, "汎用スキル 15本",
                     font_size=Pt(16), text_color=WHITE, bold=True)
    add_text_box(slide, Inches(3.7), Inches(1.35), Inches(5.0), Inches(0.5),
                 "全プロジェクト共通 — install.sh で配布",
                 font_size=Pt(13), color=SAGE, bold=True)

    universal_skills = [
        "session", "dashboard", "health-check", "skill-manage", "skill-evolve",
        "repo-manage", "backup", "commit-change", "config-audit",
        "worktree-cleanup", "update-record", "env-check", "platform-info",
        "claude-master", "validate"
    ]
    for i, name in enumerate(universal_skills):
        col = i % 5
        row = i // 5
        x = Inches(0.5 + col * 2.5)
        y = Inches(2.1 + row * 0.55)
        add_rounded_rect(slide, x, y, Inches(2.3), Inches(0.45),
                         SAGE, name, font_size=Pt(11), text_color=WHITE,
                         bold=False)

    # ドメインスキル ラベル
    add_rounded_rect(slide, Inches(0.5), Inches(4.0), Inches(3.5), Inches(0.6),
                     TERRACOTTA, "ドメインスキル 26本",
                     font_size=Pt(16), text_color=WHITE, bold=True)
    add_text_box(slide, Inches(4.2), Inches(4.05), Inches(6.0), Inches(0.5),
                 "機能安全専門 — /repo-manage sync で個別配布",
                 font_size=Pt(13), color=TERRACOTTA, bold=True)

    # ドメインスキル カテゴリ別
    categories = [
        ("プロセス管理", TERRACOTTA,
         ["execute-phase", "assess-spice", "manage-tbd", "validate",
          "problem-resolve"]),
        ("安全要件", DEEP_TERRA,
         ["select-standard", "safety-concept", "srs-generate", "fmea",
          "switch-standard"]),
        ("設計", OLIVE,
         ["system-design", "sw-design", "motor-control", "mcu-config",
          "hw-review"]),
        ("テスト/検証", RGBColor(0x6B, 0x85, 0x5B),
         ["test-design", "systest-design", "test-coverage", "safety-verify"]),
        ("コード/分析", RGBColor(0x8B, 0x6A, 0x4F),
         ["driver-gen", "static-analysis", "safety-diag"]),
        ("文書/他", RGBColor(0x9B, 0x8A, 0x6F),
         ["trace", "memory-map", "generate-docs", "ingest"]),
    ]

    y_base = Inches(4.8)
    for ci, (cat_name, color, skills) in enumerate(categories):
        col = ci % 3
        row = ci // 3
        x_base = Inches(0.5 + col * 4.3)
        y = y_base + Inches(row * 1.35)

        # カテゴリラベル
        add_text_box(slide, x_base, y, Inches(2.0), Inches(0.35),
                     cat_name, font_size=Pt(12), color=color, bold=True)

        # スキルリスト
        for si, skill in enumerate(skills):
            sx = x_base + Inches((si % 3) * 1.4)
            sy = y + Inches(0.35 + (si // 3) * 0.4)
            add_rounded_rect(slide, sx, sy, Inches(1.3), Inches(0.35),
                             color, skill, font_size=Pt(9), text_color=WHITE)

    add_bottom_bar(slide)
    add_slide_number(slide, 4)


# ============================================================
# Slide 5: 汎用スキル 15本
# ============================================================
def slide_05_universal_skills(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "汎用スキル 15本 — 全PJ共通の運用基盤",
                 font_size=Pt(28), color=WHITE, bold=True)

    groups = [
        ("セッション管理", TERRACOTTA, Inches(0.5), Inches(1.4),
         [("session", "セッション開始/終了/反省会"),
          ("backup", "バックアップ・タグ管理"),
          ("dashboard", "プロジェクト概要表示")]),
        ("品質保証", SAGE, Inches(6.8), Inches(1.4),
         [("health-check", "プロジェクト健全性監査"),
          ("skill-manage", "スキル品質監査/一覧"),
          ("skill-evolve", "PDCA改善サイクル")]),
        ("プロジェクト運用", OLIVE, Inches(0.5), Inches(4.2),
         [("repo-manage", "PJレジストリ/スキル同期"),
          ("backup", "Gitタグによるバックアップ"),
          ("commit-change", "コミットワークフロー"),
          ("update-record", "プロセス記録更新")]),
        ("環境・基盤", MUSTARD, Inches(6.8), Inches(4.2),
         [("env-check", "開発ツール検証"),
          ("platform-info", "OS/シェル情報表示"),
          ("config-audit", "設定ファイル監査"),
          ("worktree-cleanup", "WT整理/ブランチ削除"),
          ("claude-master", "技術巡回/最新化")]),
    ]

    for group_name, color, gx, gy, skills in groups:
        w = Inches(6.0)
        h = Inches(0.2 + len(skills) * 0.55 + 0.7)

        # グループ背景
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  gx, gy, w, h, fill_color=WHITE,
                  line_color=color, line_width=Pt(2))

        # グループヘッダー
        add_rounded_rect(slide, gx + Inches(0.15), gy + Inches(0.15),
                         Inches(3.5), Inches(0.5), color, group_name,
                         font_size=Pt(16), text_color=WHITE, bold=True)

        # スキルリスト
        for i, (sname, sdesc) in enumerate(skills):
            sy = gy + Inches(0.8 + i * 0.55)
            # スキル名バッジ
            add_rounded_rect(slide, gx + Inches(0.3), sy,
                             Inches(2.2), Inches(0.42), color, sname,
                             font_size=Pt(11), text_color=WHITE, bold=True)
            # 説明
            add_text_box(slide, gx + Inches(2.7), sy + Inches(0.05),
                         Inches(3.0), Inches(0.4), sdesc,
                         font_size=Pt(12), color=CHARCOAL)

    add_bottom_bar(slide)
    add_slide_number(slide, 5)


# ============================================================
# Slide 6: ドメインスキル — V-model配置
# ============================================================
def slide_06_domain_skills(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(11.0), Inches(0.7),
                 "ドメインスキル 26本 — V-modelに沿った機能安全開発",
                 font_size=Pt(28), color=WHITE, bold=True)

    # V-model: 左下がり（定義→設計→実装）
    v_left = [
        (Inches(0.3), Inches(1.5), "要件定義", TERRACOTTA,
         ["select-standard", "safety-concept", "srs-generate", "fmea",
          "switch-standard"]),
        (Inches(1.5), Inches(2.8), "設計", OLIVE,
         ["system-design", "sw-design", "mcu-config", "hw-review",
          "memory-map"]),
        (Inches(2.7), Inches(4.1), "実装", RGBColor(0x8B, 0x6A, 0x4F),
         ["driver-gen", "motor-control", "static-analysis"]),
    ]

    # V-model: 右上がり（テスト→検証→妥当性確認）
    v_right = [
        (Inches(8.0), Inches(1.5), "妥当性確認", DEEP_TERRA,
         ["safety-verify", "validate", "safety-diag"]),
        (Inches(6.8), Inches(2.8), "テスト", RGBColor(0x6B, 0x85, 0x5B),
         ["test-design", "systest-design", "test-coverage"]),
        (Inches(5.6), Inches(4.1), "プロセス管理", SAGE,
         ["execute-phase", "assess-spice", "manage-tbd", "problem-resolve"]),
    ]

    def draw_v_section(items):
        for vx, vy, label, color, skills in items:
            # ラベル
            add_rounded_rect(slide, vx, vy, Inches(2.5), Inches(0.5),
                             color, label, font_size=Pt(15), text_color=WHITE,
                             bold=True)
            # スキルチップ
            for si, skill in enumerate(skills):
                sx = vx + Inches(0.1 + (si % 3) * 1.6)
                sy = vy + Inches(0.6 + (si // 3) * 0.38)
                chip_w = Inches(1.5) if len(skill) <= 14 else Inches(1.7)
                add_rounded_rect(slide, sx, sy, chip_w, Inches(0.33),
                                 color, skill, font_size=Pt(9),
                                 text_color=WHITE)

    draw_v_section(v_left)
    draw_v_section(v_right)

    # V字の底部
    # 文書系
    doc_x = Inches(4.0)
    doc_y = Inches(5.2)
    add_rounded_rect(slide, doc_x, doc_y, Inches(2.2), Inches(0.5),
                     RGBColor(0x9B, 0x8A, 0x6F), "文書・データ",
                     font_size=Pt(14), text_color=WHITE, bold=True)
    doc_skills = ["generate-docs", "trace", "ingest"]
    for si, skill in enumerate(doc_skills):
        add_rounded_rect(slide, doc_x + Inches(si * 1.5), doc_y + Inches(0.6),
                         Inches(1.4), Inches(0.33),
                         RGBColor(0x9B, 0x8A, 0x6F), skill,
                         font_size=Pt(9), text_color=WHITE)

    # V字ライン（装飾）
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.0), Inches(5.7), Inches(4.5), Inches(0.04),
              fill_color=TERRACOTTA)
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(5.5), Inches(5.7), Inches(4.5), Inches(0.04),
              fill_color=TERRACOTTA)

    add_text_box(slide, Inches(3.5), Inches(6.0), Inches(6.0), Inches(0.4),
                 "左: 定義→設計→実装　　右: テスト→検証→妥当性確認",
                 font_size=Pt(12), color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)
    add_slide_number(slide, 6)


# ============================================================
# Slide 7: ナレッジベース（16ファイル）
# ============================================================
def slide_07_knowledge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "ナレッジベース — 16ファイル",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 汎用ナレッジ
    add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(4.0), Inches(0.6),
                     SAGE, "汎用ナレッジ（7本）",
                     font_size=Pt(18), text_color=WHITE, bold=True)
    add_text_box(slide, Inches(4.8), Inches(1.35), Inches(5.0), Inches(0.5),
                 "~/.claude/knowledge/ に配布 → 全PJで利用可能",
                 font_size=Pt(13), color=SAGE, bold=True)

    universal_k = [
        ("claude_code_ops.md", "コンテキスト管理・メモリ設計"),
        ("claude_platform_updates.md", "技術巡回履歴・プラットフォーム更新"),
        ("skill_lifecycle.md", "スキル成熟度モデル(L1-L5)定義"),
        ("skill_feedback_log.md", "スキル実行フィードバック記録"),
        ("git_worktree_branch_management.md", "Gitワークツリー・ブランチ戦略"),
        ("cross_platform_dev.md", "Windows/WSL/macOS互換性ガイド"),
        ("memory_paths.md", "メモリパス規約・slug算出ルール"),
    ]
    for i, (fname, desc) in enumerate(universal_k):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(2.1 + row * 0.5)
        add_rounded_rect(slide, x, y, Inches(2.8), Inches(0.42),
                         SAGE, fname, font_size=Pt(10), text_color=WHITE)
        add_text_box(slide, x + Inches(2.9), y + Inches(0.03),
                     Inches(3.2), Inches(0.4), desc,
                     font_size=Pt(11), color=CHARCOAL)

    # ドメインナレッジ
    add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(4.5), Inches(0.6),
                     TERRACOTTA, "ドメインナレッジ（9本）",
                     font_size=Pt(18), text_color=WHITE, bold=True)
    add_text_box(slide, Inches(5.2), Inches(4.35), Inches(5.0), Inches(0.5),
                 ".claude/knowledge/ に格納 → TORASAN管理下で利用",
                 font_size=Pt(13), color=TERRACOTTA, bold=True)

    domain_k = [
        ("iso26262_iec60730.md", "ISO 26262 / IEC 60730 規格要件"),
        ("product_standard_mapping.md", "製品カテゴリ→規格マッピング"),
        ("automotive_spice.md", "Automotive SPICE プロセスモデル"),
        ("safety_case_gsn.md", "GSN安全論証ノーテーション"),
        ("safety_diagnostics.md", "安全診断メカニズム設計"),
        ("misra_c_2012.md", "MISRA C:2012 コーディング規約"),
        ("fmea_guide.md", "FMEA故障モード影響解析ガイド"),
        ("bldc_safety.md", "BLDCモータ制御安全要件"),
        ("srs_template.md", "SRSテンプレート・要件仕様"),
    ]
    for i, (fname, desc) in enumerate(domain_k):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(5.1 + row * 0.45)
        add_rounded_rect(slide, x, y, Inches(2.8), Inches(0.38),
                         TERRACOTTA, fname, font_size=Pt(10),
                         text_color=WHITE)
        add_text_box(slide, x + Inches(2.9), y + Inches(0.01),
                     Inches(3.2), Inches(0.38), desc,
                     font_size=Pt(11), color=CHARCOAL)

    add_bottom_bar(slide)
    add_slide_number(slide, 7)


# ============================================================
# Slide 8: 配布モデル
# ============================================================
def slide_08_distribution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "配布モデル — install.sh & /repo-manage sync",
                 font_size=Pt(28), color=WHITE, bold=True)

    # TORASAN ハブ（中央上）
    add_rounded_rect(slide, Inches(4.5), Inches(1.5), Inches(4.3), Inches(1.2),
                     TERRACOTTA, bold=True, font_size=Pt(20),
                     text_color=WHITE)
    hub_shape = slide.shapes[-1]
    tf = hub_shape.text_frame
    tf.paragraphs[0].text = "TORASAN"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = ".claude/skills/ (41)  |  .claude/knowledge/ (16)"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xFF, 0xEE, 0xDD)
    p2.alignment = PP_ALIGN.CENTER

    # フロー1: install.sh → ~/.claude/ (汎用)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(2.0), Inches(0.4),
                 "install.sh", font_size=Pt(16), color=SAGE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_arrow(slide, Inches(2.0), Inches(3.6), Inches(1.5), Inches(0.4), SAGE)

    add_rounded_rect(slide, Inches(3.8), Inches(3.3), Inches(3.5), Inches(1.0),
                     SAGE, bold=True)
    box1 = slide.shapes[-1]
    tf1 = box1.text_frame
    tf1.paragraphs[0].text = "~/.claude/"
    tf1.paragraphs[0].font.size = Pt(18)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = WHITE
    p = tf1.add_paragraph()
    p.text = "skills/ (15本)  +  knowledge/ (7本)"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xEE, 0xFF, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    add_arrow(slide, Inches(7.6), Inches(3.6), Inches(1.3), Inches(0.4), SAGE)

    add_rounded_rect(slide, Inches(9.2), Inches(3.3), Inches(3.3), Inches(1.0),
                     LIGHT_SAGE, "全プロジェクト\nで自動利用",
                     font_size=Pt(14), text_color=OLIVE, bold=True,
                     line_color=SAGE)

    # フロー2: /repo-manage sync → 個別PJ (ドメイン)
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(3.0), Inches(0.4),
                 "/repo-manage sync", font_size=Pt(16), color=TERRACOTTA,
                 bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_arrow(slide, Inches(2.5), Inches(5.6), Inches(1.3), Inches(0.4),
              TERRACOTTA)

    add_rounded_rect(slide, Inches(4.0), Inches(5.2), Inches(3.3), Inches(1.0),
                     TERRACOTTA, bold=True)
    box2 = slide.shapes[-1]
    tf2 = box2.text_frame
    tf2.paragraphs[0].text = "ドメインスキル配布"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = WHITE
    p2 = tf2.add_paragraph()
    p2.text = "26本のドメインスキルを個別配布"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xFF, 0xEE, 0xDD)
    p2.alignment = PP_ALIGN.CENTER

    add_arrow(slide, Inches(7.6), Inches(5.5), Inches(1.3), Inches(0.4),
              TERRACOTTA)

    # 個別PJボックス群
    pj_names = ["PJ-A", "PJ-B", "PJ-C"]
    for i, name in enumerate(pj_names):
        y = Inches(5.0 + i * 0.5)
        add_rounded_rect(slide, Inches(9.2 + i * 0.2), y,
                         Inches(2.5), Inches(0.42),
                         LIGHT_TERRA, f"{name} (.claude/skills/)",
                         font_size=Pt(10), text_color=CHARCOAL,
                         line_color=TERRACOTTA)

    # 矢印: TORASAN → 両フロー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(6.6), Inches(2.7), Inches(0.04), Inches(0.6),
              fill_color=CHARCOAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(2.5), Inches(4.5), Inches(0.04), Inches(0.7),
              fill_color=CHARCOAL)

    add_bottom_bar(slide)
    add_slide_number(slide, 8)


# ============================================================
# Slide 9: スキル成熟度モデル
# ============================================================
def slide_09_maturity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "スキル成熟度モデル — L1 → L5",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 階段状のレベル表示
    levels = [
        ("L1", "Draft", "ファイル存在", LIGHT_SAGE),
        ("L2", "Active", "+ナレッジ参照", SAGE),
        ("L3", "Practiced", "+FB記録(1回以上)", MUSTARD),
        ("L4", "Mature", "+改善1回以上, 評価≥4.0", TERRACOTTA),
        ("L5", "Optimized", "+改善2回以上, 評価≥4.5", DEEP_TERRA),
    ]

    for i, (lvl, name, criteria, color) in enumerate(levels):
        x = Inches(0.5 + i * 1.0)
        y = Inches(5.0 - i * 0.8)
        w = Inches(12.0 - i * 1.0)
        h = Inches(0.7)

        # ステップ
        step = add_rounded_rect(slide, x, y, w, h, color)
        tf = step.text_frame
        tf.paragraphs[0].text = f"{lvl}  {name}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = FONT_EN
        p2 = tf.add_paragraph()
        p2.text = criteria
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xDD)
        p2.alignment = PP_ALIGN.LEFT

    # PDCA サイクル（右側）
    add_text_box(slide, Inches(8.0), Inches(1.3), Inches(5.0), Inches(0.5),
                 "PDCA 改善サイクル",
                 font_size=Pt(18), color=CHARCOAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    pdca_items = [
        ("Plan", "ギャップ分析", SAGE),
        ("Do", "スキル改善適用", TERRACOTTA),
        ("Check", "フィードバック収集", MUSTARD),
        ("Act", "ナレッジ更新", OLIVE),
    ]
    cx, cy = Inches(9.5), Inches(2.3)
    positions = [
        (Inches(9.2), Inches(1.9)),   # Plan - 上
        (Inches(10.7), Inches(2.7)),  # Do - 右
        (Inches(9.2), Inches(3.5)),   # Check - 下
        (Inches(7.7), Inches(2.7)),   # Act - 左
    ]

    for (px, py), (label, desc, color) in zip(positions, pdca_items):
        box = add_rounded_rect(slide, px, py, Inches(1.8), Inches(0.7), color)
        tf = box.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = FONT_EN
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xDD)
        p2.alignment = PP_ALIGN.CENTER

    # 循環矢印（テキストで表現）
    add_text_box(slide, Inches(10.5), Inches(2.2), Inches(0.5), Inches(0.5),
                 "→", font_size=Pt(20), color=CHARCOAL,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text_box(slide, Inches(10.5), Inches(3.6), Inches(0.5), Inches(0.5),
                 "←", font_size=Pt(20), color=CHARCOAL,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text_box(slide, Inches(9.7), Inches(1.5), Inches(0.5), Inches(0.5),
                 "↑", font_size=Pt(20), color=CHARCOAL,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text_box(slide, Inches(9.7), Inches(4.0), Inches(0.5), Inches(0.5),
                 "↓", font_size=Pt(20), color=CHARCOAL,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

    # 実行コマンド
    add_rounded_rect(slide, Inches(0.5), Inches(5.9), Inches(8.0), Inches(0.7),
                     WHITE, line_color=SAGE)
    cmd_shape = slide.shapes[-1]
    tf = cmd_shape.text_frame
    tf.paragraphs[0].text = \
        "/skill-evolve feedback → /skill-evolve analyze → /skill-evolve improve"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = SAGE
    tf.paragraphs[0].font.name = FONT_EN
    tf.paragraphs[0].font.bold = True

    add_bottom_bar(slide)
    add_slide_number(slide, 9)


# ============================================================
# Slide 10: まとめ & 効果
# ============================================================
def slide_10_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_SAND)

    # ヘッダー
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=CHARCOAL)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.7),
                 "まとめ — TORASAN がもたらす3つの価値",
                 font_size=Pt(32), color=WHITE, bold=True)

    # 3つの価値（大きな円）
    values = [
        ("再利用性", "一度作ったスキルを\n全PJで共有・活用",
         "install.sh で\nワンコマンド配布", TERRACOTTA, Inches(1.5)),
        ("品質向上", "成熟度モデルL1→L5\nPDCAで継続的改善",
         "フィードバック駆動の\n品質管理", SAGE, Inches(5.0)),
        ("自動化", "AI (Claude) が\nスキルを実行",
         "人手を介さず\n一貫した品質", OLIVE, Inches(8.5)),
    ]

    for title, desc1, desc2, color, cx in values:
        # メイン円
        circle = add_circle(slide, cx, Inches(1.8), Inches(3.0), color,
                            bold=True)
        tf = circle.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = ""
        p2.font.size = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc1
        p3.font.size = Pt(13)
        p3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xEE)
        p3.alignment = PP_ALIGN.CENTER

    # キーメッセージ
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5),
              fill_color=WHITE, line_color=TERRACOTTA, line_width=Pt(3))

    msg_tf = add_text_box(
        slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.2),
        "", font_size=Pt(16), color=CHARCOAL, alignment=PP_ALIGN.CENTER)
    msg_tf.paragraphs[0].text = \
        "TORASAN = スキル(実行する知恵) + ナレッジ(参照する知識) の統合管理"
    msg_tf.paragraphs[0].font.size = Pt(18)
    msg_tf.paragraphs[0].font.bold = True
    msg_tf.paragraphs[0].font.color.rgb = CHARCOAL

    p2 = msg_tf.add_paragraph()
    p2.text = "41スキル + 16ナレッジが連携し、機能安全開発の全フェーズをAIが自動実行"
    p2.font.size = Pt(15)
    p2.font.color.rgb = TERRACOTTA
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # ボトムバー (テラコッタ)
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
              fill_color=TERRACOTTA)
    add_slide_number(slide, 10)


# ============================================================
# メイン
# ============================================================
def main():
    prs = Presentation()

    # ワイドスクリーン 16:9
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 全スライド生成
    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_skill_vs_knowledge(prs)
    slide_04_skill_overview(prs)
    slide_05_universal_skills(prs)
    slide_06_domain_skills(prs)
    slide_07_knowledge(prs)
    slide_08_distribution(prs)
    slide_09_maturity(prs)
    slide_10_summary(prs)

    # 保存
    output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "TORASAN_Skills_Knowledge.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
