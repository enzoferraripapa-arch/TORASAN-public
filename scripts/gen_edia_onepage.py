"""EDIA roadmap - 元スライドの忠実な再現"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# -- 元スライドの色をそのまま使用 --
BLUE     = RGBColor(0x00, 0x00, 0xFF)    # 上段ヘッダー
DKBLUE   = RGBColor(0x00, 0x20, 0x60)    # 下段ヘッダー
BLACK    = RGBColor(0x00, 0x00, 0x00)
RED      = RGBColor(0xFF, 0x00, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW   = RGBColor(0xFF, 0xFF, 0x00)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
LTBLUE   = RGBColor(0xCC, 0xE5, 0xFF)    # 凡例: 技術的可能性
LTPINK   = RGBColor(0xFF, 0xE0, 0xE0)    # 凡例: 壁
LTGREEN  = RGBColor(0xE0, 0xFF, 0xE0)    # 凡例: 実効適用
FJ = None  # テーマフォントをそのまま使用


# -- Helpers --
def sh(stype, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(stype, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def add_shadow(shape):
    """元スライドと同じ影パラメータ"""
    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "40000", "dist": "23000",
        "dir": "5400000", "rotWithShape": "0"
    })
    srgb = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgb.makeelement(qn("a:alpha"), {"val": "38000"})
    srgb.append(alpha)
    outerShdw.append(srgb)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def add_swoosh_arrow(slide, left, top, width, height, rotation_deg,
                     color1=(0x4F, 0x81, 0xBD), color2=(0xDB, 0xE5, 0xF1)):
    """swooshArrow (曲線矢印) をグラデーション付きで追加。
    color1: 濃い側, color2: 薄い側"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    spPr = s._element.spPr

    # 回転設定
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.set("rot", str(int(rotation_deg * 60000)))

    # prstGeom を swooshArrow に変更
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "swooshArrow")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = prstGeom.makeelement(qn("a:avLst"), {})
            prstGeom.append(avLst)
        else:
            avLst.clear()
        gd1 = avLst.makeelement(qn("a:gd"), {"name": "adj1", "fmla": "val 25000"})
        gd2 = avLst.makeelement(qn("a:gd"), {"name": "adj2", "fmla": "val 25000"})
        avLst.append(gd1)
        avLst.append(gd2)

    # グラデーション塗り (XML直接構築)
    # まず既存の fill を削除
    for tag in ["a:solidFill", "a:gradFill", "a:noFill"]:
        old = spPr.find(qn(tag))
        if old is not None:
            spPr.remove(old)
    gradFill = spPr.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})
    # Stop 1: 濃い色 (0%)
    gs1 = gsLst.makeelement(qn("a:gs"), {"pos": "0"})
    clr1 = gs1.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % color1})
    alpha1 = clr1.makeelement(qn("a:alpha"), {"val": "80000"})  # 80%不透明
    clr1.append(alpha1)
    gs1.append(clr1)
    gsLst.append(gs1)
    # Stop 2: 薄い色 (100%)
    gs2 = gsLst.makeelement(qn("a:gs"), {"pos": "100000"})
    clr2 = gs2.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % color2})
    alpha2 = clr2.makeelement(qn("a:alpha"), {"val": "50000"})  # 50%不透明
    clr2.append(alpha2)
    gs2.append(clr2)
    gsLst.append(gs2)
    gradFill.append(gsLst)
    # 線形グラデーション: 右→左 (矢印先端が濃く、尾が薄い)
    lin = gradFill.makeelement(qn("a:lin"), {"ang": "10800000", "scaled": "1"})
    gradFill.append(lin)
    # spPr の prstGeom の後に挿入
    prstGeom_el = spPr.find(qn("a:prstGeom"))
    if prstGeom_el is not None:
        prstGeom_el.addnext(gradFill)
    else:
        spPr.append(gradFill)

    # 線なし
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = spPr.makeelement(qn("a:ln"), {})
        spPr.append(ln)
    else:
        ln.clear()
    noFill = ln.makeelement(qn("a:noFill"), {})
    ln.append(noFill)

    # エフェクト: glow + reflection
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    glow = effectLst.makeelement(qn("a:glow"), {"rad": "127000"})
    srgb = glow.makeelement(qn("a:srgbClr"), {"val": "4BACC6"})
    lumMod2 = srgb.makeelement(qn("a:lumMod"), {"val": "20000"})
    lumOff2 = srgb.makeelement(qn("a:lumOff"), {"val": "80000"})
    a2 = srgb.makeelement(qn("a:alpha"), {"val": "10000"})
    srgb.append(lumMod2)
    srgb.append(lumOff2)
    srgb.append(a2)
    glow.append(srgb)
    effectLst.append(glow)
    refl = effectLst.makeelement(qn("a:reflection"), {
        "stA": "0", "endPos": "65000", "dist": "50800",
        "dir": "5400000", "sy": "-100000", "algn": "bl",
        "rotWithShape": "0"
    })
    effectLst.append(refl)
    spPr.append(effectLst)

    return s


def tb(l, t, w, h, text="", sz=9, color=BLACK, bold=True, align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def para(tf, text, sz=9, color=BLACK, bold=True, sp=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = sp
    return p


def rr(l, t, w, h, fill, text="", sz=14, tc=WHITE, bold=True, shadow=True):
    s = sh(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=fill)
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = tc
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    if shadow:
        add_shadow(s)
    return s


# ============================================================
# 座標 (元スライドから抽出)
# ============================================================
# Phaseヘッダー座標 (上段: 技術分析) - 隙間を詰めて配置
UPPER_HDR = [
    (Inches(0.08), Inches(3.55)),    # Phase 1
    (Inches(3.25), Inches(2.25)),    # Phase 2
    (Inches(6.44), Inches(1.98)),    # Phase 3
    (Inches(9.61), Inches(1.67)),    # Phase 4
]
HDR_W = Inches(3.10)
HDR_H = Inches(0.55)

# Phaseラベル座標 (上段)
UPPER_LBL = [
    (Inches(0.12), Inches(3.28)),
    (Inches(3.30), Inches(1.98)),
    (Inches(6.49), Inches(1.71)),
    (Inches(9.66), Inches(1.40)),
]

# コンテンツ座標 (上段) - ヘッダーとの隙間を詰める
UPPER_BODY = [
    (Inches(0.02), Inches(4.12)),
    (Inches(3.20),  Inches(2.82)),
    (Inches(6.38),  Inches(2.55)),
    (Inches(9.56),  Inches(2.24)),
]
BODY_W = Inches(3.15)
BODY_H = Inches(1.35)

# Phaseヘッダー座標 (下段: 施策) - 隙間を詰める
LOWER_HDR = [
    None,
    (Inches(3.25), Inches(4.95)),
    (Inches(6.44), Inches(4.60)),
    (Inches(9.61), Inches(4.23)),
]

# Phaseラベル座標 (下段)
LOWER_LBL = [
    None,
    (Inches(3.30), Inches(4.68)),
    (Inches(6.49), Inches(4.33)),
    (Inches(9.66), Inches(3.96)),
]

# コンテンツ座標 (下段)
LOWER_BODY = [
    (Inches(3.20), Inches(5.52)),
    None,
    (Inches(6.38), Inches(5.17)),
    (Inches(9.56), Inches(4.80)),
]
LOWER_BODY_H = [Inches(1.35), None, Inches(1.55), Inches(2.00)]

# 縦線 (Phase間の区切り)
VLINES = [
    (Inches(3.20), Inches(1.90), Inches(5.0)),
    (Inches(6.39), Inches(1.55), Inches(5.0)),
    (Inches(9.56), Inches(1.05), Inches(5.0)),
]

# ============================================================
# 1) 背景の swooshArrow (曲線矢印 - デザインの要)
# ============================================================
# 図形47: 大きなスウッシュ (上段) - ユーザー調整済み座標
add_swoosh_arrow(slide,
                 Inches(0.0), Inches(2.41), Inches(13.0), Inches(4.0),
                 rotation_deg=355.0,
                 color1=(0x4F, 0x81, 0xBD), color2=(0xDB, 0xE5, 0xF1))

# 図形13: 大きなスウッシュ (下段) - ユーザー調整済み座標
add_swoosh_arrow(slide,
                 Inches(0.2), Inches(4.38), Inches(12.5), Inches(3.3),
                 rotation_deg=1.5,
                 color1=(0x1F, 0x49, 0x7D), color2=(0xC6, 0xD9, 0xF0))

# ============================================================
# 2) タイトル
# ============================================================
tb(Inches(0.16), Inches(0.06), Inches(13.02), Inches(0.54),
   "EDIA \u2014 AI\u00d7\u7d44\u8fbc\u307f\u958b\u767a \u30ed\u30fc\u30c9\u30de\u30c3\u30d7",
   sz=28, color=BLACK, bold=True)

# サブタイトル帯 (黄色背景 + 赤枠)
sub_box = sh(MSO_SHAPE.RECTANGLE,
             Inches(0), Inches(0.74), Inches(13.333), Inches(0.69),
             fill=YELLOW, line=RED, lw=Pt(3))
sub_tf = tb(Inches(0.1), Inches(0.76), Inches(13.1), Inches(0.32),
            "", sz=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
sub_tf.paragraphs[0].text = (
    "\u30c9\u30e1\u30a4\u30f3\u7279\u5316AI\u30a8\u30fc\u30b8\u30a7\u30f3\u30c8\u306b\u3088\u308b"
    "\u7d44\u8fbc\u307f\u958b\u767a\u30d7\u30ed\u30bb\u30b9\u306e\u81ea\u52d5\u5316 \u2014 "
    "\u6280\u8853\u7684\u53ef\u80fd\u6027\u3068\u5b9f\u88c5\u30ae\u30e3\u30c3\u30d7\u306e\u5206\u6790"
)
p2 = para(sub_tf,
          "\u300cAI\u3092\u5c0e\u5165\u3059\u308b\u304b\u5426\u304b\u300d\u3067\u306f\u306a\u304f"
          "\u300cAI\u306b\u3069\u308c\u3060\u3051\u8cea\u306e\u9ad8\u3044\u30c9\u30e1\u30a4\u30f3"
          "\u77e5\u8b58\u3092\u98df\u308f\u305b\u3089\u308c\u308b\u304b\u300d\u304c\u7af6\u4e89"
          "\u529b\u3092\u6c7a\u3081\u308b",
          sz=18, color=BLACK, bold=True)
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# 3) 凡例 (左上)
# ============================================================
# 凡例ボックス
legend = sh(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.09), Inches(1.62), Inches(3.01), Inches(0.46),
            fill=LGRAY)
legend.line.fill.background()
ltf = tb(Inches(0.12), Inches(1.62), Inches(2.95), Inches(0.15),
         "\u3010\u51e1\u4f8b\u3011", sz=9, color=BLACK, bold=True)
para(ltf, "\u25ce\u53ef\u80fd\uff1a\u6280\u8853\u7684\u306b\u5b9f\u7528\u53ef\u80fd\u306a\u9818\u57df",
     sz=9, color=BLACK)
para(ltf, "\u25b3\u9650\u754c\uff1a\u6cd5\u898f\u5236\u30fb\u502b\u7406\u30fb\u60c5\u7dd2\u7684\u8981\u56e0\u3067\u5236\u7d04",
     sz=9, color=BLACK)

# フロー図: 技術的可能性 → 壁 → 実効適用
rr(Inches(0.05), Inches(2.12), Inches(0.98), Inches(0.37),
   LTBLUE, "\u6280\u8853\u7684\u53ef\u80fd\u6027\n60\uff5e90%",
   sz=9, tc=BLACK, bold=True, shadow=False)
tb(Inches(1.04), Inches(2.21), Inches(0.22), Inches(0.20),
   "\u2192", sz=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
rr(Inches(1.25), Inches(2.12), Inches(0.87), Inches(0.37),
   LTPINK, "\u6cd5\u898f\u5236\n\u502b\u7406\u30fb\u8a8d\u8a3c\n\u60c5\u7dd2\u7684\u53d7\u5bb9",
   sz=6, tc=BLACK, bold=True, shadow=False)
tb(Inches(2.13), Inches(2.21), Inches(0.22), Inches(0.20),
   "\u2192", sz=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
rr(Inches(2.35), Inches(2.12), Inches(0.82), Inches(0.37),
   LTGREEN, "\u5b9f\u52b9\u9069\u7528\n20\uff5e50%",
   sz=9, tc=BLACK, bold=True, shadow=False)

# フロー説明
tb(Inches(0.09), Inches(2.53), Inches(3.01), Inches(0.15),
   "\u30c9\u30e1\u30a4\u30f3\u30b9\u30ad\u30eb\u306e\u6574\u5099\u3067\u58c1\u3092\u6bb5\u968e\u7684\u306b\u7a81\u7834",
   sz=9, color=BLACK, bold=True)
tb(Inches(0.09), Inches(2.68), Inches(3.01), Inches(0.15),
   "\u2192 \u975e\u5b89\u5168\u9818\u57df\u304b\u3089\u5b89\u5168\u9818\u57df\u3078\u6bb5\u968e\u7684\u306b\u5c55\u958b",
   sz=9, color=BLACK, bold=True)

# ============================================================
# 4) 縦の区切り線 (2pt)
# ============================================================
for vx, vy, vh in VLINES:
    line = slide.shapes.add_connector(1, vx, vy, vx, vy + vh)
    line.line.width = Pt(2)
    line.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)

# ============================================================
# 5) 上段: Phase ラベル + ヘッダー + コンテンツ
# ============================================================
phase_years = ["Phase 1\u3000\uff12\uff10\uff12\uff16\uff5e\uff12\uff10\uff12\uff17",
               "Phase 2\u3000\uff12\uff10\uff12\uff18\uff5e\uff12\uff10\uff13\uff10",
               "Phase 3\u3000\uff12\uff10\uff13\uff11\uff5e\uff12\uff10\uff13\uff15",
               "Phase 4\u3000\uff12\uff10\uff13\uff15\uff5e"]

upper_headers = [
    "Level 1 \u6c4e\u7528AI\u30a8\u30fc\u30b8\u30a7\u30f3\u30c8\n\u57fa\u76e4\u69cb\u7bc9\u671f",
    "Level 2 \u30c9\u30e1\u30a4\u30f3\u7279\u5316AI\n\u652f\u63f4AI\u306b\u3088\u308b\u7d44\u7e54\u5c55\u958b",
    "AI\u99c6\u52d5\u578b\u8a2d\u8a08\u306e\u5168\u793e\u6a19\u6e96\u5316\n\u30d5\u30eb\u30d1\u30a4\u30d7\u30e9\u30a4\u30f3\u69cb\u7bc9",
    "\u81ea\u5f8b\u578bAI\u8a2d\u8a08\n\u7523\u696d\u5909\u9769\u671f",
]

upper_contents = [
    ("\u25ce\u53ef\u80fd\uff1a\u6587\u66f8\u81ea\u52d5\u751f\u6210\u3001\uff7f\uff8c\uff84\uff73\uff6a\uff71\u8a2d\u8a08\uff83\uff7d\uff84\u81ea\u52d5\u5316",
     "\u3000\u3000\u3000\u3001\u69cb\u6210\u7ba1\u7406\u81ea\u52d5\u5316",
     "\u3000\u3000\u3000\u3001MISRA-C\u6e96\u62e0\uff7a\uff70\uff84\uff9e\u751f\u6210\u3001\u9759\u7684\u89e3\u6790\u81ea\u52d5\u5b9f\u884c",
     "\u25b3\u9650\u754c\uff1a\u898f\u683c\u304c\u4eba\u9593\u306b\u3088\u308b\u5b89\u5168\u5224\u65ad\u3092\u524d\u63d0",
     "\u3000\u3000\u3000\u6cd5\u7684\u8cac\u4efb\u30fb\u8a8d\u8a3c\u304cAI\u6210\u679c\u7269\u3092\u672a\u60f3\u5b9a",
     "\u3000\u3000\u3000\uff84\uff9e\uff92\uff72\uff9d\u77e5\u8b58\u4e0d\u8db3\u3067\u5b9f\u52b9\u9069\u752820-30%",
     "\u6280\u8853\u768460-70%\u53ef\u80fd \u2192 \u5b9f\u52b920-30%\u306e\uff77\uff9e\uff6c\uff6f\uff8c\uff9f"),
    ("\u25ce\u53ef\u80fd\uff1a\uff7c\uff7d\uff83\uff91\u8981\u4ef6\u5206\u6790\u304c\u8df3\u8e8d",
     "\u3000\u3000\u3000AI\u652f\u63f4\u304c\u4e00\u822c\u5316",
     "\u3000\u3000\u3000\uff8f\uff99\uff81\uff6a\uff70\uff7c\uff9e\uff6a\uff9d\uff84\u306b\u3088\u308bV\uff8c\uff9f\uff9b\uff7e\uff7d\u534a\u81ea\u52d5\u5316",
     "\u25b3\u9650\u754c\uff1a\u5b89\u5168\u6a5f\u80fd\u306f\u60c5\u7dd2\u7684\u53d7\u5bb9\u306e\u58c1",
     "\u3000\u3000\u3000\u8a8d\u8a3c\u6a5f\u95a2\u306eAI\u6210\u679c\u7269\u53d7\u5165\u57fa\u6e96\u304c\u767a\u5c55\u9014\u4e0a",
     "\u3000\u3000\u3000\uff7b\uff8c\uff9f\uff97\uff72\uff81\uff6a\uff70\uff9d\u5168\u4f53\u306eAI\u6210\u719f\u5ea6\u683c\u5dee",
     "\u6280\u8853\u768480-90%\u53ef\u80fd \u2192 \u5b9f\u52b940-50%\u3078\u6539\u5584"),
    ("\u25ce\u53ef\u80fd\uff1a\u5b89\u5168\u6a5f\u80fd\u306eAI\u652f\u63f4\u8a2d\u8a08",
     "\u3000\u3000\u3000Formal Verification\u00d7AI\u3067\u8a3c\u660e\u652f\u63f4",
     "\u3000\u3000\u3000AI\u8a2d\u8a08\u306e\u6cd5\u7684\uff8c\uff9a\uff70\uff91\uff9c\uff70\uff78\u304c\u56fd\u969b\u6a19\u6e96\u5316",
     "\u25b3\u9650\u754c\uff1aHuman-on-the-Loop\u304c\u4e3b\u6d41",
     "\u3000\u3000\u3000(\u5b8c\u5168\u81ea\u5f8b\u306f\u4e0d\u53ef)",
     "\u3000\u3000\u3000\u6700\u7d42\u5b89\u5168\u5224\u65ad\u30fb\u502b\u7406\u5224\u65ad\u306f\u4eba\u9593\u5fc5\u9808",
     "\u5b9f\u52b9\u9069\u7528\u7bc4\u56f2\u306e\u5927\u5e45\u62e1\u5927",
     "\u8a8d\u8a3c\uff8c\uff9f\uff9b\uff7e\uff7d\u306eAI\u5bfe\u5fdc\u5b8c\u4e86"),
    ("\u25ce\u53ef\u80fd\uff1aAI\u304c\u8a2d\u8a08\u4e3b\u4f53\u3001\u4eba\u9593\u304c\u76e3\u7763\u8005",
     "\u3000\u3000\u3000\u91cf\u5b50\uff7a\uff9d\uff8b\uff9f\uff6d\uff70\uff83\uff68\uff9d\uff78\uff9e\u9023\u643a\u306e\u5f62\u5f0f\u691c\u8a3c\u9ad8\u901f\u5316",
     "\u3000\u3000\u3000\u696d\u754c\u5168\u4f53\u306e\u958b\u767a\u751f\u7523\u6027\u304c\u6841\u9055\u3044\u306b\u5411\u4e0a",
     "\u25b3\u6b8b\u5b58\u8ab2\u984c\uff1a\u5b89\u5168\uff78\uff98\uff83\uff68\uff76\uff99\u3067\u306f\u4eba\u9593\u6700\u7d42\u627f\u8a8d\u6b8b\u5b58",
     "\u3000\u3000\u3000AI\u8d77\u56e0\u91cd\u5927\u4e8b\u6545\u6642\u306e\u793e\u4f1a\u7684\u5f8c\u9000\uff98\uff7d\uff78",
     "\u3000\u3000\u3000\uff8c\uff99\uff75\uff70\uff84\u8a2d\u8a08\u306e\u6cd5\u7684\u5b9a\u7fa9\u3068\u8cac\u4efb\u4f53\u7cfb\u78ba\u7acb",
     "\u300c\u5f79\u5272\u9006\u8ee2\u300d\u304c\u6280\u8853\u7684\u306b\u6210\u7acb"),
]

# Phase labels (上段)
for i in range(4):
    lx, ly = UPPER_LBL[i]
    tb(lx, ly, Inches(2.81), Inches(0.27),
       phase_years[i], sz=16, color=BLACK, bold=True)

# Phase headers (上段 - 青)
for i in range(4):
    hx, hy = UPPER_HDR[i]
    rr(hx, hy, HDR_W, HDR_H, BLUE, upper_headers[i], sz=14, bold=True)

# Phase contents (上段 - 枠なしテキスト)
for i in range(4):
    bx, by = UPPER_BODY[i]
    bh = Inches(1.46) if i >= 2 else BODY_H
    lines = upper_contents[i]
    tf = tb(bx, by, BODY_W, bh, "", sz=9, color=BLACK, bold=True)
    for j, line in enumerate(lines):
        is_summary = (j == len(lines) - 1) or (j == len(lines) - 2 and i >= 2)
        if line.startswith("\u25ce"):
            c = BLUE
        elif line.startswith("\u25b3"):
            c = RED
        elif is_summary and not line.startswith("\u3000"):
            c = BLUE
        else:
            c = BLACK
        if j == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = c
            tf.paragraphs[0].font.bold = True
        else:
            para(tf, line, sz=9, color=c, bold=True)

# ============================================================
# 6) 下段: Phase ラベル + ヘッダー + コンテンツ
# ============================================================
lower_headers = [
    None,
    "\u30c9\u30e1\u30a4\u30f3\u30b9\u30ad\u30eb\u6574\u5099\u3068\nAI\u30ac\u30d0\u30ca\u30f3\u30b9\u57fa\u76e4\u69cb\u7bc9",
    "\u5b89\u5168\u9818\u57df\u3078\u306e\u6bb5\u968e\u5c55\u958b\n\u8a8d\u8a3c\u5b9f\u7e3e\u69cb\u7bc9",
    "\u5168\u88fd\u54c1\u30e9\u30a4\u30f3\u5c55\u958b\n\u30b5\u30d7\u30e9\u30a4\u30c1\u30a7\u30fc\u30f3AI\u5316",
]

lower_contents = [
    # Phase 1 施策 (x位置はPhase2列)
    ["\u25c7\u6700\u512a\u5148\u65bd\u7b56",
     "\u25c6\u5168\u793eAI\uff76\uff9e\uff8a\uff9e\uff85\uff9d\uff7d\uff8c\uff9a\uff70\uff91\uff9c\uff70\uff78\u7b56\u5b9a",
     "\u25c6\uff84\uff9e\uff92\uff72\uff9d\uff7d\uff77\uff99\u306e\u5168\u793e\u6574\u5099(6\u9818\u57df)",
     "\u25c6\u5148\u884c\uff8a\uff9f\uff72\uff9b\uff6f\uff84PJ 3-5\u4ef6\u3067\u5b9a\u91cfROI\u53d6\u5f97",
     "\u25c7\u4eba\u6750\u80b2\u6210",
     "\u25c6AI\uff7d\uff77\uff99\uff8f\uff7d\uff80\uff70\u4eba\u6750\u306e\u80b2\u6210\u958b\u59cb",
     "\u25c6NDA\u30fb\u6a5f\u5bc6\u6027\u7ba1\u7406\uff8c\uff9a\uff70\uff91\uff9c\uff70\uff78\u7b56\u5b9a",
     "\u25c6\u300c\uff84\uff9e\uff92\uff72\uff9d\uff7d\uff77\uff99\u3092\u66f8\u3051\u308b\u4eba\u6750\u300d\u306e\u6226\u7565\u7684\u80b2\u6210"],
    # Phase 2 施策 → rendered at LOWER_BODY[2] position
    None,
    ["\u25c7\u6280\u8853\u5c55\u958b",
     "\u25c6\u5b89\u5168\u9818\u57df\u3078\u306e\u6bb5\u968e\u7684AI\u9069\u7528",
     "\u25c6\u8a8d\u8a3c\u6a5f\u95a2\u3068\u306e\u5bfe\u8a71\u30fbAI\u6210\u679c\u7269\u53d7\u5165\u5b9f\u7e3e\u69cb\u7bc9",
     "\u25c6\uff8f\uff99\uff81\uff6a\uff70\uff7c\uff9e\uff6a\uff9d\uff84V\uff8c\uff9f\uff9b\uff7e\uff7d\u306e\u5b9f\u7528\u5316",
     "\u25c7\u7d44\u7e54\u30fb\uff7b\uff8c\uff9f\uff97\uff72\uff81\uff6a\uff70\uff9d",
     "\u25c6\uff7b\uff8c\uff9f\uff97\uff72\uff94\u5411\u3051AI\u6d3b\u7528\u8981\u4ef6\u306e\u8abf\u9054\u4ed5\u69d8\u5316",
     "\u25c6\u5168\u793e\u6a2a\u65ad\u306e\uff8d\uff9e\uff7d\uff84\uff8c\uff9f\uff97\uff78\uff83\uff68\uff7d\u5c55\u958b",
     "\u25c6\u696d\u754c\u6a19\u6e96\u5316\u6d3b\u52d5\u3078\u306e\uff98\uff70\uff80\uff9e\uff70\uff7c\uff6f\uff8c\uff9f\u767a\u63ee",
     "\u2192\uff76\uff9d\uff8a\uff9f\uff86\uff70\u500b\u5225\u304b\u3089\u5168\u793e\u7d71\u4e00AI\u6d3b\u7528\u3078"],
    ["\u25c7\u5168\u793e\u5c55\u958b",
     "\u25c6\u5168\u88fd\u54c1\uff97\uff72\uff9d\u3067\u306eAI\u652f\u63f4\u8a2d\u8a08\u6a19\u6e96\u5316",
     "\u25c6\uff7b\uff8c\uff9f\uff97\uff72\uff81\uff6a\uff70\uff9d\u5168\u4f53\u306eAI\u5316\u63a8\u9032",
     "\u25c6\u6b21\u4e16\u4ee3\uff71\uff70\uff77\uff83\uff78\uff81\uff6c\u306e\u691c\u8a0e\u958b\u59cb",
     "\u25c7\u6301\u7d9a\u7684\u7af6\u4e89\u529b",
     "\u25c6\u5dee\u5225\u5316\u6280\u8853(\u542bAI)\u306e\u6700\u5927\u9650\u6d3b\u7528",
     "\u25c6\u77ed\u671f\u9593\u3067\u591a\u304f\u306eSMC\u3092\u7af6\u5408\u306b\u5148\u99c6\u3051\u4e0a\u5e02",
     "\u25c6\uff78\uff9e\uff9b\uff70\uff8a\uff9e\uff99\u898f\u5236\u74b0\u5883\u3067\u306e\uff98\uff70\uff80\uff9e\uff70\uff7c\uff6f\uff8c\uff9f",
     "\u2192\u7523\u696d\uff6a\uff7a\uff7c\uff7d\uff83\uff91\u306e\u518d\u5b9a\u7fa9\u30fb\u65b0\uff8b\uff9e\uff7c\uff9e\uff88\uff7d\uff93\uff83\uff9e\uff99"],
]

# Phase labels (下段)
for i in range(1, 4):
    lx, ly = LOWER_LBL[i]
    yr = phase_years[i]
    lbl_w = Inches(3.25) if i == 1 else Inches(2.45)
    lbl_h = Inches(0.37) if i == 1 else Inches(0.27)
    tb(lx, ly, lbl_w, lbl_h,
       yr, sz=16, color=BLACK, bold=True)

# Phase headers (下段 - ダークブルー)
for i in range(1, 4):
    hx, hy = LOWER_HDR[i]
    rr(hx, hy, HDR_W, HDR_H, DKBLUE, lower_headers[i], sz=14, bold=True)

# Phase contents (下段 - 枠なしテキスト)
content_positions = [
    (Inches(3.20), Inches(5.52), Inches(1.35)),   # Phase 1 施策
    None,
    (Inches(6.38), Inches(5.17), Inches(1.55)),   # Phase 3 施策
    (Inches(9.56), Inches(4.80), Inches(2.00)),   # Phase 4 施策
]

for i in [0, 2, 3]:
    bx, by, bh = content_positions[i]
    lines = lower_contents[i]
    tf = tb(bx, by, BODY_W, bh, "", sz=9, color=BLACK, bold=True)
    for j, line in enumerate(lines):
        if line.startswith("\u25c7"):
            c = DKBLUE
        elif line.startswith("\u2192"):
            c = DKBLUE
        else:
            c = BLACK
        if j == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = c
            tf.paragraphs[0].font.bold = True
        else:
            para(tf, line, sz=9, color=c, bold=True)

# ============================================================
# 7) 左下: ギャップ分析 (吹き出し + テキスト)
# ============================================================
# 吹き出し背景 (wedgeRectCallout, 回転なし, ポインタ上向き)
callout = sh(MSO_SHAPE.RECTANGULAR_CALLOUT,
             Inches(0.12), Inches(5.98), Inches(2.94), Inches(0.92),
             fill=DKBLUE)
# prstGeom を wedgeRectCallout に変更 + adj でポインタ上方向
prstGeom = callout._element.spPr.find(qn("a:prstGeom"))
if prstGeom is not None:
    prstGeom.set("prst", "wedgeRectCallout")
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = prstGeom.makeelement(qn("a:avLst"), {})
        prstGeom.append(avLst)
    else:
        avLst.clear()
    gd1 = avLst.makeelement(qn("a:gd"), {"name": "adj1", "fmla": "val 60180"})
    gd2 = avLst.makeelement(qn("a:gd"), {"name": "adj2", "fmla": "val -95818"})
    avLst.append(gd1)
    avLst.append(gd2)
add_shadow(callout)

# ギャップ要因 (元スライド: Shape 30)
tf_gap = tb(Inches(0.10), Inches(5.61), Inches(1.90), Inches(1.68),
            "ギャップ要因", sz=9, color=WHITE, bold=True)
for item in ["◇法的責任", "◇倫理的判断",
             "◇情緒的受容", "◇認証・規格"]:
    para(tf_gap, item, sz=9, color=WHITE, bold=True)

# 解決の方向性 (元スライド: Shape 31)
tf_sol = tb(Inches(1.78), Inches(5.61), Inches(1.53), Inches(1.68),
            "解決の方向性", sz=9, color=WHITE, bold=True)
for item in ["◇段階的導入", "◇判断根拠可視化",
             "◇規制形成参画"]:
    para(tf_sol, item, sz=9, color=WHITE, bold=True)

# ============================================================
# 8) 右下: 凡例バッジ
# ============================================================
rr(Inches(11.88), Inches(6.43), Inches(1.42), Inches(0.35),
   BLUE, "\u30c9\u30e1\u30a4\u30f3\u7279\u5316AI (Level 2)", sz=9, bold=True)
rr(Inches(11.88), Inches(6.83), Inches(1.42), Inches(0.35),
   DKBLUE, "\uff84\uff9e\uff92\uff72\uff9d\uff7d\uff77\uff99\u306b\u3088\u308b\u7d44\u7e54\u77e5\u306e\u5f62\u5f0f\u77e5\u5316",
   sz=9, bold=True)

# ============================================================
# Save
# ============================================================
out = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "samples",
                   "EDIA_roadmap_onepage.pptx")
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"Generated: {out}")
